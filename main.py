import os
import random
import string
import pyautogui
from datetime import date
from pynput.keyboard import Key, Listener
from pptx import Presentation
from pptx.util import Inches
from pydrive.drive import GoogleDrive 
from pydrive.auth import GoogleAuth


def on_press(key):
    global image_num
    if key == Key.insert:
        #add screenshot
        myScreenshot = pyautogui.screenshot()
        image_name = day_path + "\\screenshots\\image_" + str(image_num) +".png"
        myScreenshot = pyautogui.screenshot()
        myScreenshot.save(image_name)
        print(image_name)
        print('{0} pressed'.format(
        key))
        image_num += 1

    elif key == Key.delete:
        #delete image
        if image_num > 1:
            if not os.path.isfile(day_path + "\\deleted_images\\former_image" + str(image_num - 1) + ".png"):
                os.rename(day_path + "\\screenshots\\image_" + str(image_num - 1) +".png",
                              day_path + "\\deleted_images\\former_image" + str(image_num - 1) + ".png")
            else:
                random_letters = ''.join(random.choice(string.ascii_lowercase) for i in range(4))
                os.rename(day_path + "\\screenshots\\image_" + str(image_num - 1) +".png",
                              day_path + "\\deleted_images\\former_image" + str(image_num - 1)+ random_letters + ".png")
                
            image_num -= 1
            print('{0} pressed'.format(key))
            
    if key == Key.end:
        # make powerpoint, add powerpoint to Google Drive, and close
        make_powerpoint()
        upload_to_drive()
        return False

    if key == Key.esc:
        return False

def make_master_dir():
    if not os.path.isdir(day_path):
        os.mkdir(day_path)
        os.mkdir(day_path + "\\screenshots")
        os.mkdir(day_path + "\\deleted_images")
        print("paths made")
    else:
        print("already made")


def make_powerpoint():
    #make powerpoint
    global d4
    global lec_num
    
    powerpoint_folder = "C:\\Users\\lpc\\Desktop\\chem\\slide_decks"
    pic_folder = "C:\\Users\\lpc\\Desktop\\chem\\" + d4 + "\\screenshots"
    if not os.path.isdir(powerpoint_folder):
        os.mkdir(powerpoint_folder)
    prs = Presentation()
    num = 6
    layout = prs.slide_layouts[int(num)]
    for pic in os.listdir(pic_folder):
        slide = prs.slides.add_slide(layout)
        
        height = Inches(5)
        width = Inches(10)
        finished_slide = slide.shapes.add_picture(os.path.join(pic_folder,pic), 0, Inches(0.9), width=width)
        
    lec_num = input("Input lecture number:")
    prs.save(powerpoint_folder + "\\Lecture_" +
                str(lec_num) + "_" + d4 + ".pptx")
    os.startfile(powerpoint_folder + "\\Lecture_" +
                str(lec_num) + "_" + d4 + ".pptx")
    print("File saved as: " + powerpoint_folder + "\\Lecture_" +
                str(lec_num) + "_" + d4 + ".pptx")

    
def upload_to_drive():
    #log into Google account and upload powerpoint
    gauth = GoogleAuth()
    gauth.LocalWebserverAuth()        
    drive = GoogleDrive(gauth)
    x = "Lecture_" + str(lec_num) + "_Dec-25-2020.pptx"
    print(x)
    path = "C:\\Users\\lpc\\Desktop\\chem\\slide_decks"
    f = drive.CreateFile({'title': x}) 
    f.SetContentFile(os.path.join(path, x)) 
    f.Upload()
    print("File uploaded as: " + os.path.join(path, x))

        
master = "C:\\Users\\lpc\\Desktop\\chem\\"
today = date.today()
d4 = today.strftime("%b-%d-%Y")
day_path = master + d4
image_num = 0

make_master_dir()


with Listener(
        on_press=on_press) as listener:
    listener.join()
